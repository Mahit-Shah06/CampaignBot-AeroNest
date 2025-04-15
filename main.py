# campaignbot_backend/main.py - Enhanced FastAPI Backend with GPT-4

from fastapi import FastAPI, Request
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import openai
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import tempfile

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

openai.api_key = "YOUR_OPENAI_API_KEY"

class CampaignData(BaseModel):
    brand_name: str
    objective: str
    audience: str
    competitors: str
    visual_style: str

@app.post("/generate-campaign")
def generate_campaign(data: CampaignData):
    prompt = f"""
    You are an expert digital marketing strategist AI. Based on the information below, generate a full campaign plan including:
    1. Unique Value Proposition
    2. Suggested Slogan
    3. Ad Formats by Platform (Instagram, Google, YouTube, etc.)
    4. Ad Copy (Instagram caption, Google ad headline, CTA button text, Email subject line)
    5. Recommended Target Audience breakdown
    6. Ideal KPIs (click-through, conversions, etc.)
    7. Estimated Budget Plan (for small, medium, and large campaigns)

    Details:
    Brand Name: {data.brand_name}
    Objective: {data.objective}
    Target Audience: {data.audience}
    Competitors: {data.competitors}
    Visual Style: {data.visual_style}
    """

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a professional ad agency strategist."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.85,
            max_tokens=1000
        )
        return {"summary": response['choices'][0]['message']['content']}
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.get("/generate-logo")
def generate_logo(prompt: str):
    try:
        response = openai.Image.create(
            prompt=prompt,
            n=1,
            size="512x512"
        )
        return {"url": response['data'][0]['url']}
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)

@app.post("/export-docx")
def export_docx(data: CampaignData):
    doc = Document()
    doc.add_heading("Campaign Plan", 0)
    doc.add_paragraph(f"Brand Name: {data.brand_name}")
    doc.add_paragraph(f"Objective: {data.objective}")
    doc.add_paragraph(f"Audience: {data.audience}")
    doc.add_paragraph(f"Competitors: {data.competitors}")
    doc.add_paragraph(f"Visual Style: {data.visual_style}")
    doc.add_paragraph("Deliverables: Ad Plan, Visual Identity, A/B Testing Strategy")
    temp = tempfile.NamedTemporaryFile(delete=False, suffix=".docx")
    doc.save(temp.name)
    return FileResponse(temp.name, filename="campaign_plan.docx")

@app.post("/export-pptx")
def export_pptx(data: CampaignData):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Campaign Plan"
    content.text = (
        f"Brand Name: {data.brand_name}\n"
        f"Objective: {data.objective}\n"
        f"Audience: {data.audience}\n"
        f"Competitors: {data.competitors}\n"
        f"Visual Style: {data.visual_style}\n"
        "Deliverables: Ad Plan, Visual Identity, A/B Testing Strategy"
    )
    temp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(temp.name)
    return FileResponse(temp.name, filename="campaign_plan.pptx")
